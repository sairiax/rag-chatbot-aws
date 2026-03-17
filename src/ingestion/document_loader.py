"""
Document loader supporting PDF, Word, PowerPoint, plain text, and images.
OCR is applied automatically to image-only pages or image files.
"""
from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING, List, Optional

from langchain_core.documents import Document
from loguru import logger

if TYPE_CHECKING:
    from src.ingestion.ocr_processor import OCRProcessor

# Maps file extension → private loader method name
_EXT_TO_METHOD: dict[str, str] = {
    ".pdf": "_load_pdf",
    ".docx": "_load_docx",
    ".doc": "_load_docx",
    ".pptx": "_load_pptx",
    ".ppt": "_load_pptx",
    ".txt": "_load_text",
    ".md": "_load_text",
    ".png": "_load_image",
    ".jpg": "_load_image",
    ".jpeg": "_load_image",
    ".tiff": "_load_image",
    ".bmp": "_load_image",
}

SUPPORTED_EXTENSIONS: list[str] = list(_EXT_TO_METHOD.keys())

# Minimum characters per page to consider as "text present" (avoids OCR on digital PDFs)
_MIN_TEXT_CHARS = 50


class MultiFormatDocumentLoader:
    """Load documents from multiple file formats into LangChain Documents."""

    def __init__(self, ocr_processor: Optional["OCRProcessor"] = None) -> None:
        self.ocr_processor = ocr_processor

    # ──────────────────────────────────────────────────────────────────────
    # Public API
    # ──────────────────────────────────────────────────────────────────────

    def load_file(self, file_path: str | Path) -> List[Document]:
        """Load a single file. Returns a list of Documents with rich metadata."""
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = path.suffix.lower()
        if ext not in _EXT_TO_METHOD:
            raise ValueError(
                f"Unsupported format '{ext}'. Supported: {SUPPORTED_EXTENSIONS}"
            )

        loader_method = getattr(self, _EXT_TO_METHOD[ext])
        documents: List[Document] = loader_method(path)

        # Attach common metadata to every chunk
        for doc in documents:
            doc.metadata.setdefault("source", str(path))
            doc.metadata.setdefault("filename", path.name)
            doc.metadata.setdefault("file_type", ext.lstrip("."))

        logger.info(f"Loaded {len(documents)} document(s) from '{path.name}'")
        return documents

    def load_directory(
        self, dir_path: str | Path, recursive: bool = True
    ) -> List[Document]:
        """Load all supported files inside a directory."""
        path = Path(dir_path)
        pattern = "**/*" if recursive else "*"

        all_documents: List[Document] = []
        for ext in _EXT_TO_METHOD:
            for file_path in path.glob(f"{pattern}{ext}"):
                try:
                    all_documents.extend(self.load_file(file_path))
                except Exception as exc:
                    logger.error(f"Failed to load '{file_path.name}': {exc}")

        logger.info(f"Directory scan complete — {len(all_documents)} document(s) loaded")
        return all_documents

    # ──────────────────────────────────────────────────────────────────────
    # Private loaders
    # ──────────────────────────────────────────────────────────────────────

    def _load_pdf(self, path: Path) -> List[Document]:
        """Load PDF. Falls back to OCR for pages with little or no text."""
        from langchain_community.document_loaders import PyPDFLoader

        loader = PyPDFLoader(str(path))
        documents = loader.load()

        if self.ocr_processor:
            processed: List[Document] = []
            for doc in documents:
                page_num = doc.metadata.get("page", 0)
                if len(doc.page_content.strip()) < _MIN_TEXT_CHARS:
                    logger.info(
                        f"Sparse text on page {page_num + 1} of '{path.name}' — applying OCR"
                    )
                    ocr_text = self.ocr_processor.extract_text_from_pdf_page(
                        str(path), page_num
                    )
                    if ocr_text:
                        doc.page_content = ocr_text
                        doc.metadata["ocr_applied"] = True
                processed.append(doc)
            return processed

        return documents

    def _load_docx(self, path: Path) -> List[Document]:
        """Load .docx / .doc Word document."""
        from langchain_community.document_loaders import Docx2txtLoader

        return Docx2txtLoader(str(path)).load()

    def _load_pptx(self, path: Path) -> List[Document]:
        """Load .pptx / .ppt PowerPoint — one Document per slide."""
        from pptx import Presentation

        prs = Presentation(str(path))
        documents: List[Document] = []
        total_slides = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, start=1):
            parts: List[str] = []

            # Slide title
            title = ""
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
                parts.append(f"Título: {title}")

            # All text shapes (skip duplicating the title)
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and text != title:
                        parts.append(text)

            # Speaker notes
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                if notes_tf and notes_tf.text.strip():
                    parts.append(f"Notas: {notes_tf.text.strip()}")

            content = "\n".join(parts)
            if content.strip():
                documents.append(
                    Document(
                        page_content=content,
                        metadata={
                            "slide_number": slide_num,
                            "slide_title": title,
                            "total_slides": total_slides,
                        },
                    )
                )

        return documents

    def _load_text(self, path: Path) -> List[Document]:
        """Load plain text or Markdown file."""
        from langchain_community.document_loaders import TextLoader

        return TextLoader(str(path), encoding="utf-8").load()

    def _load_image(self, path: Path) -> List[Document]:
        """Load image via OCR. Requires OCR processor to be configured."""
        if not self.ocr_processor:
            raise ValueError(
                "An OCRProcessor is required to load image files. "
                "Pass one to MultiFormatDocumentLoader(ocr_processor=...)."
            )
        text = self.ocr_processor.extract_text_from_image(str(path))
        return [Document(page_content=text, metadata={"ocr_applied": True})]
